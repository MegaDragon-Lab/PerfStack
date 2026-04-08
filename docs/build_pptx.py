"""Generate PerfStack.pptx presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0f, 0x11, 0x17)   # near-black
CARD      = RGBColor(0x1a, 0x1f, 0x2e)   # dark-blue card
ACCENT    = RGBColor(0x63, 0x66, 0xf1)   # indigo
ACCENT2   = RGBColor(0x81, 0x8c, 0xf8)   # lighter indigo
GREEN     = RGBColor(0x22, 0xc5, 0x5e)
YELLOW    = RGBColor(0xf5, 0x9e, 0x0b)
PURPLE    = RGBColor(0xa7, 0x8b, 0xfa)
WHITE     = RGBColor(0xf1, 0xf5, 0xf9)
MUTED     = RGBColor(0x94, 0xa3, 0xb8)
BORDER    = RGBColor(0x2d, 0x37, 0x48)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

DOCS = os.path.dirname(os.path.abspath(__file__))
ARCH_PNG  = os.path.join(DOCS, "architecture.png")
INFRA_PNG = os.path.join(DOCS, "infra.png")


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=CARD, border_color=BORDER, border_pt=1):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, x, y, w, h, size=14, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def accent_bar(slide, y=Inches(0.08)):
    bar = slide.shapes.add_shape(1, 0, y, W, Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


# ── Slide 1 — Title ───────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank_slide(prs)
    bg(s)

    # gradient-ish wide accent strip
    strip = s.shapes.add_shape(1, 0, Inches(2.8), W, Inches(2.1))
    strip.fill.solid()
    strip.fill.fore_color.rgb = CARD
    strip.line.fill.background()

    accent_bar(s, Inches(2.8))
    accent_bar(s, Inches(4.88))

    label(s, "PerfStack", Inches(1), Inches(3.0), Inches(11), Inches(1.3),
          size=60, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    label(s, "Load Testing Platform", Inches(1), Inches(4.1), Inches(11), Inches(0.7),
          size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    label(s, "K6  ·  Grafana  ·  InfluxDB  ·  Kubernetes",
          Inches(1), Inches(4.75), Inches(11), Inches(0.5),
          size=16, color=MUTED, align=PP_ALIGN.CENTER)

    label(s, "FICO — epc_owner@fico.com", Inches(1), Inches(6.8), Inches(11), Inches(0.4),
          size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 2 — Agenda ─────────────────────────────────────────────────────────
def slide_agenda(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Agenda", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)

    items = [
        ("01", "The Problem",          "Why load testing matters before going to production"),
        ("02", "Architecture",         "Components, data flow, and deployment topology"),
        ("03", "Tech Stack",           "K6, FastAPI, React, InfluxDB, Grafana, Kubernetes"),
        ("04", "How it Works",         "End-to-end flow: from UI click to rendered report"),
        ("05", "Reporting",            "Metrics, KPIs, Grafana charts — all in one HTML report"),
        ("06", "Live Demo",            "Running a real load test against a target API"),
    ]

    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        bx = Inches(0.5) + col * Inches(6.4)
        by = Inches(1.3) + row * Inches(1.8)
        bw = Inches(6.0)
        bh = Inches(1.55)

        b = box(s, bx, by, bw, bh)
        label(s, num, bx + Inches(0.18), by + Inches(0.1), Inches(0.6), Inches(0.55),
              size=22, bold=True, color=ACCENT)
        label(s, title, bx + Inches(0.7), by + Inches(0.1), bw - Inches(0.8), Inches(0.45),
              size=18, bold=True, color=WHITE)
        label(s, desc, bx + Inches(0.18), by + Inches(0.65), bw - Inches(0.3), Inches(0.75),
              size=12, color=MUTED)


# ── Slide 3 — The Problem ─────────────────────────────────────────────────────
def slide_problem(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "The Problem", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)
    label(s, "Performance issues discovered in production — not before.",
          Inches(0.6), Inches(0.95), Inches(12), Inches(0.45),
          size=18, color=MUTED)

    pains = [
        ("No visibility",       "APIs fail under real concurrency.\nProblems only surface after deployment."),
        ("Manual & slow",       "Running k6 locally means copy-pasting\ncommands, reading raw JSON output."),
        ("No shared context",   "Individual developers test in silos.\nNo historical comparison, no team reports."),
        ("OAuth complexity",    "Target APIs require bearer tokens.\nToken lifecycle is a manual chore."),
    ]

    for i, (title, body) in enumerate(pains):
        col = i % 2
        row = i // 2
        bx = Inches(0.5) + col * Inches(6.4)
        by = Inches(1.7) + row * Inches(2.5)
        b = box(s, bx, by, Inches(6.0), Inches(2.2))
        label(s, "✕", bx + Inches(0.18), by + Inches(0.12), Inches(0.5), Inches(0.5),
              size=22, bold=True, color=RGBColor(0xef, 0x44, 0x44))
        label(s, title, bx + Inches(0.65), by + Inches(0.1), Inches(5.1), Inches(0.45),
              size=17, bold=True, color=WHITE)
        label(s, body, bx + Inches(0.18), by + Inches(0.65), Inches(5.5), Inches(1.4),
              size=13, color=MUTED)

    label(s, "PerfStack solves all of this — one platform, zero manual steps.",
          Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
          size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


# ── Slide 4 — Architecture ────────────────────────────────────────────────────
def slide_architecture(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Architecture", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)

    if os.path.exists(ARCH_PNG):
        s.shapes.add_picture(ARCH_PNG, Inches(0.4), Inches(1.05),
                             Inches(12.5), Inches(6.2))
    else:
        label(s, "[architecture.png not found]", Inches(1), Inches(3), Inches(10), Inches(1),
              size=18, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 5 — Tech Stack ──────────────────────────────────────────────────────
def slide_techstack(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Tech Stack", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)

    stack = [
        ("K6",          GREEN,  "Load testing engine — VU-based, JS scripts,\nStreams metrics to InfluxDB via xk6-output-influxdb"),
        ("k6 Operator", GREEN,  "Kubernetes CRD — runs K6 as parallelized Jobs\n(4 pods × configured VUs)"),
        ("FastAPI",     ACCENT, "Python backend — IAM auth, job lifecycle,\nGrafana render API, HTML report generation"),
        ("React + Vite",ACCENT, "Frontend SPA — test config, live pod monitor,\nprogressive chart loader"),
        ("InfluxDB v2", YELLOW, "Time-series store — receives K6 metrics in real-time\nQueried via InfluxQL compatibility layer"),
        ("Grafana 12",  YELLOW, "Dashboard + image renderer — panels rendered\nto PNG via headless Chromium"),
        ("K3d",         PURPLE, "Local Kubernetes cluster on Docker —\nproduction-identical topology on a laptop"),
        ("nginx Ingress",PURPLE,"Single entry point — routes /api → backend,\n/ → frontend, /grafana → Grafana"),
    ]

    for i, (name, color, desc) in enumerate(stack):
        col = i % 2
        row = i // 2
        bx = Inches(0.4) + col * Inches(6.45)
        by = Inches(1.2) + row * Inches(1.45)
        b = box(s, bx, by, Inches(6.1), Inches(1.3))
        dot = s.shapes.add_shape(9, bx + Inches(0.18), by + Inches(0.38),
                                  Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        label(s, name, bx + Inches(0.5), by + Inches(0.08), Inches(5.3), Inches(0.42),
              size=16, bold=True, color=color)
        label(s, desc, bx + Inches(0.18), by + Inches(0.55), Inches(5.7), Inches(0.68),
              size=11, color=MUTED)


# ── Slide 6 — How it Works ────────────────────────────────────────────────────
def slide_flow(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "How it Works", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)

    steps = [
        ("①", ACCENT,  "Configure",    "Pick target URL, payload,\nVUs, duration, IAM creds\nin the React UI"),
        ("②", GREEN,   "Auth",         "Backend calls IAM OAuth2\n(client credentials flow)\nand caches Bearer token"),
        ("③", YELLOW,  "Deploy",       "K6 Operator creates a\nTestRun CRD — spins up\n4 parallel K6 pods"),
        ("④", PURPLE,  "Load",         "K6 pods hammer the target\nAPI and stream metrics to\nInfluxDB every second"),
        ("⑤", ACCENT2, "Report",       "Backend renders 5 Grafana\npanels progressively and\nserves self-contained HTML"),
    ]

    arrow_color = BORDER
    step_w = Inches(2.3)
    step_h = Inches(4.8)
    gap    = Inches(0.22)
    start_x = Inches(0.3)
    top_y  = Inches(1.4)

    for i, (num, color, title, body) in enumerate(steps):
        bx = start_x + i * (step_w + gap)
        b = box(s, bx, top_y, step_w, step_h)

        # number circle
        circ = s.shapes.add_shape(9, bx + Inches(0.9), top_y + Inches(0.2),
                                   Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        label(s, num, bx + Inches(0.88), top_y + Inches(0.18), Inches(0.54), Inches(0.5),
              size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)

        label(s, title, bx + Inches(0.15), top_y + Inches(0.88), step_w - Inches(0.3), Inches(0.5),
              size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        label(s, body, bx + Inches(0.15), top_y + Inches(1.5), step_w - Inches(0.3), Inches(3.0),
              size=12, color=MUTED, align=PP_ALIGN.CENTER)

        # arrow between boxes
        if i < len(steps) - 1:
            ax = bx + step_w + Inches(0.03)
            ay = top_y + step_h / 2 - Pt(6)
            label(s, "▶", ax, ay, gap, Inches(0.35),
                  size=12, color=BORDER, align=PP_ALIGN.CENTER)


# ── Slide 7 — Reporting ───────────────────────────────────────────────────────
def slide_reporting(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Performance Report", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)
    label(s, "Auto-generated HTML — no manual post-processing",
          Inches(0.6), Inches(0.95), Inches(12), Inches(0.45),
          size=16, color=MUTED)

    kpis = [
        ("Total Requests",  ACCENT,  "All HTTP calls made\nduring the test"),
        ("Peak req/s",       PURPLE,  "Max throughput observed\nfrom InfluxDB timeseries"),
        ("Error Rate",       GREEN,   "% of failed requests\nwith color thresholds"),
        ("p95 / p99",        YELLOW,  "Response time percentiles\nfrom K6 summary JSON"),
        ("Checks Passed",    GREEN,   "Custom K6 check results\npass/fail breakdown"),
    ]

    for i, (name, color, desc) in enumerate(kpis):
        bx = Inches(0.4) + i * Inches(2.55)
        by = Inches(1.6)
        b = box(s, bx, by, Inches(2.35), Inches(1.5))
        label(s, "●", bx + Inches(0.15), by + Inches(0.12), Inches(0.4), Inches(0.4),
              size=18, color=color)
        label(s, name, bx + Inches(0.5), by + Inches(0.1), Inches(1.8), Inches(0.45),
              size=13, bold=True, color=WHITE)
        label(s, desc, bx + Inches(0.15), by + Inches(0.65), Inches(2.1), Inches(0.75),
              size=11, color=MUTED)

    charts = [
        ("Virtual Users",            "Live VU ramp-up / ramp-down"),
        ("Requests per Second",      "Throughput timeline"),
        ("Errors per Second",        "Error spikes visualised"),
        ("Response Time (p90/p95)",  "Latency percentile bands"),
        ("Response Time Heatmap",    "Distribution over time"),
    ]

    label(s, "Grafana Panels — rendered progressively as they load",
          Inches(0.6), Inches(3.3), Inches(12), Inches(0.4),
          size=15, bold=True, color=ACCENT2)

    for i, (name, desc) in enumerate(charts):
        col = i % 3
        row = i // 3
        bx = Inches(0.4) + col * Inches(4.3)
        by = Inches(3.85) + row * Inches(1.5)
        b = box(s, bx, by, Inches(4.1), Inches(1.3))
        label(s, name, bx + Inches(0.18), by + Inches(0.1), Inches(3.7), Inches(0.45),
              size=14, bold=True, color=WHITE)
        label(s, desc, bx + Inches(0.18), by + Inches(0.6), Inches(3.7), Inches(0.6),
              size=12, color=MUTED)


# ── Slide 8 — Live Demo ───────────────────────────────────────────────────────
def slide_demo(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Live Demo", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)

    steps = [
        "Open  http://localhost  in the browser",
        "Add a service — enter IAM credentials, target URL, payload",
        "Configure VUs & duration (e.g. 10 VUs · 60 s)",
        "Click  Run Test  — watch pod status in real time",
        "Click  Download Report  — see charts load progressively",
        "Open Grafana → K6 Load Testing Results for live dashboard",
    ]

    for i, step in enumerate(steps):
        bx = Inches(1.5)
        by = Inches(1.4) + i * Inches(0.82)
        num_box = box(s, Inches(0.5), by + Inches(0.04), Inches(0.7), Inches(0.55),
                      fill_color=ACCENT, border_color=None)
        label(s, str(i + 1), Inches(0.5), by + Inches(0.04), Inches(0.7), Inches(0.55),
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(s, step, bx, by, Inches(10.8), Inches(0.6),
              size=16, color=WHITE)

    label(s, "Access:  http://localhost   |   Grafana: http://localhost/grafana  (admin / admin)",
          Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
          size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 9 — Key Takeaways ───────────────────────────────────────────────────
def slide_takeaways(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Key Takeaways", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)

    points = [
        (GREEN,  "Zero manual steps",       "From IAM auth to rendered report — fully automated.\nNo YAML editing, no command line, no copy-paste."),
        (ACCENT, "Production-identical",    "K3d cluster mirrors EC2 deployment exactly.\nSame images, same manifests, same network paths."),
        (YELLOW, "Observable by default",   "Every test streams metrics to InfluxDB and Grafana\nin real time — no post-processing needed."),
        (PURPLE, "Extensible",              "Add scenarios, save services, run against any\nOAuth2-protected API without code changes."),
    ]

    for i, (color, title, body) in enumerate(points):
        bx = Inches(0.5) + (i % 2) * Inches(6.4)
        by = Inches(1.3) + (i // 2) * Inches(2.6)
        b = box(s, bx, by, Inches(6.1), Inches(2.35))
        bar = s.shapes.add_shape(1, bx, by, Inches(0.08), Inches(2.35))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        label(s, title, bx + Inches(0.25), by + Inches(0.18), Inches(5.6), Inches(0.5),
              size=18, bold=True, color=color)
        label(s, body, bx + Inches(0.25), by + Inches(0.78), Inches(5.6), Inches(1.4),
              size=13, color=MUTED)


# ── Slide 10 — Infra View ─────────────────────────────────────────────────────
def slide_infra(prs):
    s = blank_slide(prs)
    bg(s)
    accent_bar(s)

    label(s, "Infrastructure View", Inches(0.6), Inches(0.25), Inches(10), Inches(0.6),
          size=32, bold=True, color=WHITE)
    label(s, "Kubernetes resources · namespaces · CPU / RAM limits · network paths · storage",
          Inches(0.6), Inches(0.9), Inches(12.5), Inches(0.4),
          size=14, color=MUTED)

    if os.path.exists(INFRA_PNG):
        s.shapes.add_picture(INFRA_PNG, Inches(0.3), Inches(1.4),
                             Inches(12.7), Inches(5.9))
    else:
        label(s, "[infra.png not found — run: mmdc -i docs/infra.mmd -o docs/infra.png]",
              Inches(1), Inches(3.5), Inches(11), Inches(1),
              size=14, color=MUTED, align=PP_ALIGN.CENTER)


# ── Slide 11 — Q&A / Thank you ────────────────────────────────────────────────
def slide_qa(prs):
    s = blank_slide(prs)
    bg(s)

    strip = s.shapes.add_shape(1, 0, Inches(2.8), W, Inches(2.1))
    strip.fill.solid(); strip.fill.fore_color.rgb = CARD
    strip.line.fill.background()
    accent_bar(s, Inches(2.8))
    accent_bar(s, Inches(4.88))

    label(s, "Questions?", Inches(1), Inches(3.0), Inches(11), Inches(1.1),
          size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s, "github · epc_owner@fico.com",
          Inches(1), Inches(4.1), Inches(11), Inches(0.55),
          size=20, color=MUTED, align=PP_ALIGN.CENTER)
    label(s, "PerfStack — Load Testing Platform",
          Inches(1), Inches(4.72), Inches(11), Inches(0.45),
          size=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ── Build ─────────────────────────────────────────────────────────────────────
def build():
    prs = new_prs()
    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_techstack(prs)
    slide_flow(prs)
    slide_reporting(prs)
    slide_demo(prs)
    slide_takeaways(prs)
    slide_infra(prs)
    slide_qa(prs)

    out = os.path.join(DOCS, "PerfStack.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
